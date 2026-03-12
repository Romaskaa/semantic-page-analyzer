import os
from pypdf import PdfReader
import docx
from pptx import Presentation
import pandas as pd
from PIL import Image
import pytesseract

UPLOAD_DIR = "storage/uploads"

def save_file(file):

    os.makedirs(UPLOAD_DIR, exist_ok=True)

    path = os.path.join(UPLOAD_DIR, file.filename)

    with open(path, "wb") as f:
        f.write(file.file.read())

    return path

def read_txt(path):

    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def read_pdf(path):

    text = []

    reader = PdfReader(path)

    for page in reader.pages:

        page_text = page.extract_text()

        if page_text:
            text.append(page_text)

    return "\n".join(text)


def read_docx(path):

    document = docx.Document(path)

    text = []

    for p in document.paragraphs:
        text.append(p.text)

    return "\n".join(text)


def read_pptx(path):

    prs = Presentation(path)

    text = []

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                text.append(shape.text)

    return "\n".join(text)


def read_xlsx(path):

    df = pd.read_excel(path)

    return df.to_string()


def read_image(path):

    img = Image.open(path)

    text = pytesseract.image_to_string(img)

    return text

def extract_text_from_file(path):

    ext = path.lower().split(".")[-1]

    try:

        if ext == "txt":
            return read_txt(path)

        elif ext == "pdf":
            return read_pdf(path)

        elif ext == "docx":
            return read_docx(path)

        elif ext == "pptx":
            return read_pptx(path)

        elif ext in ["xls", "xlsx"]:
            return read_xlsx(path)

        elif ext in ["png", "jpg", "jpeg"]:
            return read_image(path)

        else:
            return ""

    except Exception as e:

        print("Ошибка чтения файла:", path, e)

        return ""

def load_knowledge_base():

    texts = []

    if not os.path.exists(UPLOAD_DIR):
        return ""

    files = os.listdir(UPLOAD_DIR)

    for file in files:

        path = os.path.join(UPLOAD_DIR, file)

        if os.path.isfile(path):

            text = extract_text_from_file(path)

            if text:
                texts.append(text)

    return "\n\n".join(texts)